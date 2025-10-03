import io
import os
import platform
import urllib.parse
from io import BytesIO
from typing import List, Optional, Union

import pandas as pd
import requests
from docx import Document as DocxDocument
from langchain_community.document_loaders.parsers.pdf import PyPDFParser
from langchain_core.document_loaders import Blob
from langchain_core.document_loaders.base import BaseLoader
from langchain_core.documents.base import Document
from pptx import Presentation


class SharePointClient:
    def __init__(self, site_url, site_id, access_token, obo_access_token):
        self.site_url = site_url
        self.site_id = site_id
        self.access_token = access_token
        self.obo_access_token = obo_access_token

    def get_one_drive_id(self):
        """
        This function retrieves the OneDrive ID of the current user using the Microsoft Graph API.

        Returns:
        str: The OneDrive ID of the current user.
        """
        url = "https://graph.microsoft.com/v1.0/me/drive"
        headers = {"Authorization": f"Bearer {self.obo_access_token}"}
        resp = requests.get(url, headers=headers)
        if resp.status_code == 200:
            return resp.json().get("id")
        else:
            print(f"Error fetching OneDrive ID: {resp.status_code} - {resp.text}")
            return None

    def get_drive_id(self, library_name="Documents"):
        """Resolve a document library (drive) by its display name.

        Args:
            library_name (str): Display name of the library (e.g., 'Documents', 'Shared Documents', 'Policies').
        Returns:
            str: The ID of the drive on the SharePoint site
        """

        url = f"https://graph.microsoft.com/v1.0/sites/{self.site_id}/drives"
        headers = {"Authorization": f"Bearer {self.access_token}"}
        resp = requests.get(url, headers=headers)
        if resp.status_code != 200:
            print(f"Error listing drives: {resp.status_code} - {resp.text}")
            return None
        for drive in resp.json().get("value", []):
            if drive.get("name", "").lower() == library_name.lower():
                return drive.get("id")

        return None

    def get_folder_id(self, drive_id, folder_path):
        """
        Retrieve the ID of a folder in a drive by its path.
        """
        if not drive_id:
            return None

        if not folder_path or folder_path == "":
            url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root"
        else:
            encoded_path = urllib.parse.quote(folder_path)
            url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root:/{encoded_path}"

        headers = {"Authorization": f"Bearer {self.access_token}"}
        resp = requests.get(url, headers=headers)
        if resp.status_code == 200:
            return resp.json().get("id")
        else:
            print(
                f"Error fetching folder '{folder_path}': {resp.status_code} - {resp.text}"
            )
            return None

    def get_files(self, drive_id, file_name):
        """
        This function retrieves the ID of a specified file in a drive on a SharePoint site.

        Parameters:
        drive_id (str): The ID of the drive where the file is located.
        file_name (str): The name of the file whose ID is to be retrieved.

        Returns:
        str: The ID of the specified file.
        """
        files_url = f"https://graph.microsoft.com/v1.0/sites/{self.site_id}/drives/{drive_id}/root/children"
        response = requests.get(
            files_url, headers={"Authorization": f"Bearer {self.access_token}"}
        )
        items_data = response.json()

        for item in items_data["value"]:
            if item["name"] == file_name:
                return item["id"]
        return None

    def get_all_files_in_drive(self, drive_id, file_name=".", top=5):
        """
        Get all files from a SharePoint document library (drive), including files in all folders.
        Args:
            drive_id (str): The ID of the document library (drive).
            file_name (str): The name of the file to search for. Defaults to '.' which matches all files.
            top (int): Number of recent files to retrieve (default: 5)

        Returns:
            list: List of file metadata dictionaries.
        """
        url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/search(q='{file_name}')?$orderby=lastModifiedDateTime desc&$top={max(top*2, 20)}"
        headers = {
            "Authorization": f"Bearer {self.obo_access_token}",
            "Content-Type": "application/json",
        }
        files = []
        while url:
            resp = requests.get(url, headers=headers)
            if resp.status_code == 200:
                data = resp.json()
                # Only include items that are files (not folders)
                files.extend(
                    [item for item in data.get("value", []) if "folder" not in item]
                )
                url = data.get("@odata.nextLink")  # Handle pagination
            else:
                print(f"Error: {resp.status_code} - {resp.text}")
                break
        return files

    def get_recent_onedrive_files(self, file_name=".", top=5):
        """
        Get recent files from OneDrive for the current user.

        Args:
            file_name (str): The name of the file to search for. Defaults to '.' which matches all files.
            top (int): Number of recent files to retrieve (default: 5)

        Returns:
            list: List of recent files with metadata
        """
        url = f"https://graph.microsoft.com/v1.0/me/drive/root/search(q='{file_name}')?$orderby=lastModifiedDateTime desc&$top={max(top*2, 20)}"
        headers = {
            "Authorization": f"Bearer {self.obo_access_token}",
            "Content-Type": "application/json",
        }

        response = requests.get(url, headers=headers)
        if response.status_code == 200:
            items = response.json().get("value", [])
            # Filter out folders
            files = [item for item in items if "folder" not in item]
            # Return only up to 'top' files
            return files[:top]
        else:
            print(f"Error: {response.status_code} - {response.text}")
            return []

    def copyfile(self, file_id, drive_id, folder_id):
        """
        Copy a file in OneDrive to a specified destination folder.

        Args:
            file_id (str): The ID of the file to be copied.
            drive_id (str): The ID of the destination document library.
            folder_id (str): The ID of the destination folder within the document library.
        """

        url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/copy?@microsoft.graph.conflictBehavior=replace"

        headers = {
            "Authorization": f"Bearer {self.obo_access_token}",
            "Content-Type": "application/json",
        }

        body = {"parentReference": {"driveId": drive_id, "id": folder_id}}

        response = requests.post(url, headers=headers, json=body)

        if response.status_code == 202:
            print("File copy initiated successfully.")
        else:
            print(f"Error: {response.status_code} - {response.text}")

    def get_file_download_url(self, drive_id, file_id):
        """
        This function retrieves the download URL of a specified file from a SharePoint site.

        Parameters:
        drive_id (str): The ID of the drive on the SharePoint site.
        file_id (str): The ID of the file to be downloaded.

        Returns:
        str: The download URL of the file if successful, None otherwise.
        """
        try:
            # Get the file details
            file_url = (
                f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{file_id}"
            )
            headers = {"Authorization": f"Bearer {self.obo_access_token}"}
            response = requests.get(file_url, headers=headers)
            file_data = response.json()

            # Get the download URL and file name
            download_url = file_data["@microsoft.graph.downloadUrl"]
            file_name = file_data["name"]

            return download_url

        except requests.exceptions.RequestException as e:
            print(f"Error downloading file: {file_name} err: {e}")
            return None

    def get_site_analytics(self):
        """
        Retrieve site usage analytics for the SharePoint site.

        Returns:
            dict: A dictionary containing site usage analytics data.
        """
        url = f"https://graph.microsoft.com/v1.0/sites/{self.site_id}/analytics/allTime"
        headers = {
            "Authorization": f"Bearer {self.access_token}",
            "Content-Type": "application/json",
        }
        resp = requests.get(url, headers=headers)
        if resp.status_code == 200:
            return resp.json()
        else:
            print(f"Error fetching site analytics: {resp.status_code} - {resp.text}")
            return None

    def update_file_metadata(self, drive_id, file_id, metadata: dict):
        """
        Update metadata of a file in SharePoint document library.

        Args:
            drive_id (str): The ID of the document library (drive).
            file_id (str): The ID of the file to update.
            metadata (dict): A dictionary of metadata fields to update.
        """
        headers = {
            "Authorization": f"Bearer {self.obo_access_token}",
            "Content-Type": "application/json",
        }

        # 1. Get listItemId
        url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{file_id}?expand=listItem"
        r = requests.get(url, headers=headers)
        r.raise_for_status()
        list_item_id = r.json()["listItem"]["id"]

        # 2. Get listId
        url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/list"
        r = requests.get(url, headers=headers)
        r.raise_for_status()
        list_id = r.json()["id"]

        # 3. Patch metadata
        url = f"https://graph.microsoft.com/v1.0/sites/{self.site_id}/lists/{list_id}/items/{list_item_id}/fields"
        r = requests.patch(url, headers=headers, json=metadata)

        if r.status_code == 200:
            print("Metadata updated successfully")
            return r.json()
        else:
            print(f"Error updating metadata: {r.status_code} - {r.text}")
            return None

    def load_sharepoint_document_by_name(self, drive_id, file_name):
        """
        This function retrieves a document from a SharePoint site and loads it into memory using a custom loader based on the file type.

        Args:
            drive_id (str): The ID of the drive on the SharePoint site.
            file_name (str): The name of the file to be loaded.

        Returns:
            object: A custom loader object that can handle the content of the loaded file. The type of the loader depends on the file type.
        """
        # Get the download URL and the file name by querying the Microsoft Graph API
        file_url = (
            f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root:/{file_name}"
        )
        headers = {
            "Authorization": f"Bearer {self.access_token}"
        }  # Use the stored access token for authorization
        response = requests.get(
            file_url, headers=headers
        )  # Make the HTTP request to get file details
        file_data = response.json()  # Parse the JSON response to get file data

        file_type = file_data.get("file", {}).get(
            "mimeType", ""
        )  # Get the MIME type of the file
        download_url = file_data[
            "@microsoft.graph.downloadUrl"
        ]  # Extract the direct download URL from the response

        # Get the file content from the download URL
        response = requests.get(
            download_url, headers=headers
        )  # Make the HTTP request to download the file

        # Create a BytesIO object from the response content, which allows for reading and writing bytes in memory
        stream = io.BytesIO(
            response.content
        )  # This is useful for handling binary data like files without saving to disk

        # Check the file type and use the appropriate custom loader to handle the file content
        if file_type == "application/pdf":
            # Use CustomPDFLoader to handle PDF files; it initializes with the stream and file name
            loader = CustomPDFLoader(stream, file_name)
            return loader
        elif (
            file_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            # Use CustomWordLoader for Word documents to handle and potentially split the document's content
            loader = CustomWordLoader(stream, file_name)
            return loader
        elif (
            file_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            # Use CustomPPTLoader for PowerPoint presentations to read and split the presentation into slides
            loader = CustomPPTLoader(stream, file_name)
            return loader
        elif (
            file_type
            == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ):
            # Use CustomExcelLoader for Excel spreadsheets to read and possibly split the sheets into separate parts
            loader = CustomExcelLoader(stream, file_name)
            return loader
        elif file_type in ["text/csv", "text/plain"]:
            # Use CustomTextLoader for plain text or CSV files to handle and split text as needed
            loader = CustomTextLoader(stream, file_name)
            return loader
        else:
            print(f"Unsupported file type: {file_type}")
            pass  # Placeholder for additional file types that may need to be implemented in the future


class CustomPDFLoader(BaseLoader):
    """
    This class is a custom loader for PDF files. It inherits from the BaseLoader class.

    The class is initialized with a binary stream of the PDF file, the file name, an optional password for protected PDFs,
    and a flag indicating whether to extract images from the PDF.

    The load method converts the binary stream into a Blob object, parses the PDF, and converts each page or segment
    into a separate document object. The file name is added as metadata to each document.
    """

    def __init__(
        self,
        stream: BytesIO,
        filename: str,
        password: Optional[Union[str, bytes]] = None,
        extract_images: bool = False,
    ):
        # Initialize with a binary stream, file name, optional password, and an image extraction flag
        self.stream = stream
        self.filename = filename
        # Initialize a PDF parser with optional password protection and image extraction settings
        self.parser = PyPDFParser(password=password, extract_images=extract_images)

    def load(self) -> List[Document]:
        # Convert the binary stream into a Blob object which is required by the parser
        blob = Blob.from_data(self.stream.getvalue())
        # Parse the PDF and convert each page or segment into a separate document object
        documents = list(self.parser.parse(blob))

        # Add the filename as metadata to each document for identification
        for doc in documents:
            doc.metadata.update({"source": self.filename})

        return documents


class CustomWordLoader(BaseLoader):
    """
    This class is a custom loader for Word documents. It extends the BaseLoader class and overrides its methods.
    It uses the python-docx library to parse Word documents and optionally splits the text into manageable documents.

    Attributes:
    stream (io.BytesIO): A binary stream of the Word document.
    filename (str): The name of the Word document.
    """

    def __init__(self, stream, filename: str):
        # Initialize with a binary stream and filename
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Use python-docx to parse the Word document from the binary stream
        doc = DocxDocument(self.stream)
        # Extract and concatenate all paragraph texts into a single string
        text = "\n".join([p.text for p in doc.paragraphs])

        # Check if a text splitter utility is provided
        if text_splitter is not None:
            # Use the provided splitter to divide the text into manageable documents
            split_text = text_splitter.create_documents([text])
        else:
            # Without a splitter, treat the entire text as one document
            split_text = [{"text": text, "metadata": {"source": self.filename}}]

        # Add source metadata to each resulting document
        for doc in split_text:
            if isinstance(doc, dict):
                doc["metadata"] = {**doc.get("metadata", {}), "source": self.filename}
            else:
                doc.metadata = {**doc.metadata, "source": self.filename}

        return split_text


class CustomExcelLoader(BaseLoader):
    """
    This class is a custom loader for Excel files. It inherits from the BaseLoader class.

    The class takes a binary stream of an Excel file and a filename as input, and provides a method to load the Excel file into memory and split its content into separate documents based on the sheets in the workbook.

    Attributes:
    stream (io.BytesIO): A binary stream of the Excel file.
    filename (str): The name of the Excel file.
    """

    def __init__(self, stream, filename: str):
        # Initialize with a binary stream and filename
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Use pandas to load the Excel file from the binary stream
        xls = pd.ExcelFile(self.stream, engine="openpyxl")
        # Get the list of all sheet names in the workbook
        sheet_names = xls.sheet_names

        split_sheets = []
        for sheet in sheet_names:
            # Parse each sheet into a DataFrame
            df = xls.parse(sheet)
            # Convert the DataFrame to a single string with each cell value separated by new lines
            text = "\n".join(df.values.astype(str).flatten().tolist())

            # Check if a text splitter is provided to further divide the sheet content
            if text_splitter is not None:
                # Use the splitter to create documents from the text
                split_text = text_splitter.create_documents([text])
                # Add metadata to each document
                for doc in split_text:
                    doc.metadata = {"source": self.filename, "page": sheet}
                split_sheets.extend(split_text)
            else:
                # Without a splitter, treat the entire sheet text as one document
                doc = Document(text, metadata={"source": self.filename, "page": sheet})
                split_sheets.append(doc)

        return split_sheets


class CustomPPTLoader(BaseLoader):
    """
    This class is a custom loader for PowerPoint files. It inherits from the BaseLoader class.

    The class takes a binary stream of a PowerPoint file and a filename as input, and provides a method to load the PowerPoint file into memory and split its content into separate documents based on the slides in the presentation.

    Attributes:
    stream (io.BytesIO): A binary stream of the PowerPoint file.
    filename (str): The name of the PowerPoint file.
    """

    def __init__(self, stream, filename):
        # Initialize with a binary stream and filename
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Use python-pptx to parse the PowerPoint file from the binary stream
        prs = Presentation(self.stream)
        # Prepare to collect documents
        documents = []

        # Iterate over each slide in the presentation
        for i, slide in enumerate(prs.slides):
            # Extract all text content from each slide
            slide_text = "\n".join(
                [
                    paragraph.text
                    for shape in slide.shapes
                    if shape.has_text_frame
                    for paragraph in shape.text_frame.paragraphs
                ]
            )

            # Check if a text splitter is provided
            if text_splitter is None:
                # Treat each slide's text as a single document
                doc = {
                    "text": slide_text,
                    "metadata": {"source": self.filename, "page": i + 1},
                }
                # doc = {'text': slide_text, 'metadata': {'source': self.filename}}
                documents.append(doc)
            else:
                # Use the splitter to divide the slide text into smaller documents
                split_text = text_splitter.create_documents([slide_text])
                # Add metadata and collect each document
                for doc in split_text:
                    doc.metadata = {"source": self.filename, "page": i + 1}
                    # doc.metadata = {'source': self.filename}
                documents.extend(split_text)

        return documents


import chardet


class CustomTextLoader(BaseLoader):
    """
    This class is a custom loader for text files. It inherits from the BaseLoader class.

    The class takes a binary stream of a text file and a filename as input, and provides a method to load the text file into memory and split its content into separate documents.

    Attributes:
    stream (io.BytesIO): A binary stream of the text file.
    filename (str): The name of the text file.
    """

    def __init__(self, stream, filename: str):
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Use chardet to detect the encoding of the stream
        rawdata = self.stream.read()
        result = chardet.detect(rawdata)
        text = rawdata.decode(result["encoding"])

        if text_splitter is not None:
            split_text = text_splitter.create_documents([text])
        else:
            split_text = [{"text": text, "metadata": {"source": self.filename}}]

        for doc in split_text:
            if isinstance(doc, dict):
                doc["metadata"] = {**doc.get("metadata", {}), "source": self.filename}
            else:
                doc.metadata = {**doc.metadata, "source": self.filename}

        return split_text
